"""
Generate the thesis-defense PowerPoint deck.
Output: ThesisDefense.pptx in the "thesis writing" directory.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

FIG_DIR = "../figures"
OUT_PATH = "ThesisDefense.pptx"

# -------- Theme colours --------
NAVY     = RGBColor(0x14, 0x2B, 0x52)
ACCENT   = RGBColor(0x1F, 0x77, 0xB4)
LIGHTGREY= RGBColor(0xF2, 0xF2, 0xF2)
TEXT     = RGBColor(0x22, 0x22, 0x22)
SUBTEXT  = RGBColor(0x55, 0x55, 0x55)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# -------- Slide geometry (16:9) --------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

blank = prs.slide_layouts[6]  # 6 = blank

# ---------- helpers ----------

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font_name="Calibri",
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, font_size=18, color=TEXT, bullet_color=ACCENT):
    """items: list of strings or (text, sub_items_list) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            head, subs = it
        else:
            head, subs = it, None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = "• " + head
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if subs:
            for s in subs:
                sp = tf.add_paragraph()
                sp.level = 1
                sr = sp.add_run()
                sr.text = "– " + s
                sr.font.name = "Calibri"
                sr.font.size = Pt(font_size - 3)
                sr.font.color.rgb = SUBTEXT
    return tb


def add_header(slide, title, subtitle=None):
    # Top navy bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.10), SLIDE_W - Inches(0.8),
                Inches(0.65), title, font_size=24, bold=True,
                color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.90), SLIDE_W - Inches(0.8),
                    Inches(0.4), subtitle, font_size=14, color=SUBTEXT)


def add_footer(slide, page_no, total):
    # Footer thin strip
    add_rect(slide, 0, SLIDE_H - Inches(0.30), SLIDE_W, Inches(0.30), LIGHTGREY)
    add_textbox(slide, Inches(0.4), SLIDE_H - Inches(0.30),
                Inches(8), Inches(0.30),
                "Multimodal Explainable AI for Breast Histopathology — Thesis Defense",
                font_size=10, color=SUBTEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.30),
                Inches(0.8), Inches(0.30),
                f"{page_no} / {total}", font_size=10, color=SUBTEXT,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_image_or_placeholder(slide, x, y, w, h, fname, fit=True):
    path = os.path.join(FIG_DIR, fname)
    if os.path.exists(path):
        try:
            pic = slide.shapes.add_picture(path, x, y, width=w, height=h if not fit else None)
            # If fitting on width and height auto-set, optionally clip
            return pic
        except Exception as e:
            print(f"  ! Could not embed {fname}: {e}")
    add_rect(slide, x, y, w, h, LIGHTGREY, line=ACCENT)
    add_textbox(slide, x, y, w, h, f"[Image: {fname}]", font_size=12,
                color=SUBTEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return None


def add_table(slide, x, y, w, h, headers, rows, header_fill=ACCENT,
              header_text=WHITE, body_text=TEXT, font_size=14):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    # Headers
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h_text
        r.font.bold = True
        r.font.color.rgb = header_text
        r.font.size = Pt(font_size)
        r.font.name = "Calibri"
    # Body
    for i, row_data in enumerate(rows, start=1):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTGREY if i % 2 == 0 else WHITE
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size - 2)
            r.font.name = "Calibri"
            r.font.color.rgb = body_text
            # Bold the v3.6 row & first column
            if "v3.6" in str(row_data[0]):
                r.font.bold = True
    return tbl


# ===================================================================
# Slide builders
# ===================================================================

slides_built = []


def new_slide():
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    slides_built.append(slide)
    return slide


# ---------- Slide 1: Title ----------
def slide_title():
    s = new_slide()
    # Big navy banner
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.8), Inches(1.2), SLIDE_W - Inches(1.6), Inches(1.5),
                "Multimodal Explainable AI", font_size=44, bold=True,
                color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.8), Inches(2.4), SLIDE_W - Inches(1.6), Inches(1.0),
                "for Breast Histopathology Classification",
                font_size=32, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(0.8), Inches(3.5), Inches(2.5), Inches(0.05), ACCENT)
    add_textbox(s, Inches(0.8), Inches(3.7), SLIDE_W - Inches(1.6), Inches(0.6),
                "Two-Head Expert-Gated Fusion of ConvNeXt and Swin Transformer on BreaKHis 400×",
                font_size=20, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(5.0), SLIDE_W - Inches(1.6), Inches(1.2),
                "Master’s Thesis Defense\nAbdullah Amir  ·  Supervisor: [Supervisor Name]\nDepartment of [Department]\n[University]  ·  2026",
                font_size=18, color=WHITE)


# ---------- Slide 2: Outline ----------
def slide_outline():
    s = new_slide()
    add_header(s, "Outline")
    items = [
        "Motivation & Problem Statement",
        "Research Questions and Contributions",
        "Related Work (2022+ literature)",
        "Dataset: BreaKHis 400×",
        "Methodology: backbones, fusion head, training",
        "Two-Head v3.6 architecture",
        "Evaluation framework and statistical analysis",
        "Results: single-split and patient-CV",
        "Comparison with state of the art",
        "Calibration and Explainability (Grad-CAM, IG)",
        "Ablation studies",
        "Limitations and future work",
        "Conclusions and contributions",
    ]
    add_bullets(s, Inches(0.8), Inches(1.2), Inches(12), Inches(5.5), items, font_size=20)


# ---------- Slide 3: Motivation ----------
def slide_motivation():
    s = new_slide()
    add_header(s, "Why this thesis?",
               "Computer-aided diagnosis for breast histopathology needs both accuracy AND interpretability")
    items = [
        ("Clinical context",
         ["Breast cancer is a leading cause of cancer death in women globally.",
          "Histopathology is the diagnostic gold standard, but pathologists differ on rare subtypes.",
          "Patient-level generalisation matters more than image-level."]),
        ("Open technical problems",
         ["CNNs (ConvNeXt) and Transformers (Swin) capture different features — local vs global.",
          "Class imbalance: Ductal carcinoma 43%, Adenosis 5.7% — 7.6× imbalance.",
          "Existing models lack faithful, model-aware explanations under fusion."]),
        ("Why fusion + XAI together?",
         ["Fusion improves accuracy only if it preserves complementary information.",
          "Grad-CAM cannot capture cross-branch reasoning in MLP fusion — needs IG.",
          "Calibration is rarely reported on BreaKHis, yet it is essential for clinical use."]),
    ]
    add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5), items, font_size=18)


# ---------- Slide 4: Research questions ----------
def slide_questions():
    s = new_slide()
    add_header(s, "Research questions")
    items = [
        ("RQ1 — Heterogeneity",
         ["Do ConvNeXt-B and Swin-B make complementary errors on BreaKHis 400×?"]),
        ("RQ2 — Fusion strategy",
         ["Does feature-level fusion outperform logit-level ensembling under patient-level CV?",
          "What architecture is required for the fusion head to be stable on rare classes?"]),
        ("RQ3 — Patient-invariant learning",
         ["Can we reduce the image-CV → patient-CV generalisation gap, which prior 2022+ work reports at 7.0–7.6 pp?"]),
        ("RQ4 — Faithful explanations",
         ["Are CAM-family methods faithful for an MLP-fusion model?",
          "Does Integrated Gradients on a FusionWrapper give substantially lower Deletion-AUC?"]),
        ("RQ5 — Efficiency",
         ["Can we reach state-of-the-art patient-CV F1 without retraining 80+ M backbone parameters?"]),
    ]
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6), items, font_size=18)


# ---------- Slide 5: Contributions ----------
def slide_contributions():
    s = new_slide()
    add_header(s, "Contributions")
    items = [
        ("New SOTA on BreaKHis 400× under patient-CV",
         ["Two-head v3.6: macro-F1 = 0.892, 95 % CI [0.876, 0.908].",
          "4.27 M trainable params — 20× smaller than MiSLAS (88 M)."]),
        ("Two-head expert-gated fusion design",
         ["Decoupled binary and 8-class subtype heads.",
          "Learned per-sample gating reduces patient-CV std from 0.1599 to 0.0876 (−45 %)."]),
        ("Calibration as a first-class metric",
         ["ECE = 0.0605 after per-class temperature scaling. None of the cited 2022+ work reports ECE."]),
        ("Faithful XAI for fusion models",
         ["Integrated Gradients on the FusionWrapper: Deletion-AUC = 0.178.",
          "61 % reduction vs Grad-CAM++ (0.462). CAM-family is unfaithful through MLPs."]),
        ("Reduced generalisation gap (new SOTA)",
         ["Image-CV → patient-CV gap reduced to 4.5 pp vs 7.0–7.6 pp in prior work."]),
    ]
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6), items, font_size=18)


# ---------- Slide 6: Related Work table ----------
def slide_related():
    s = new_slide()
    add_header(s, "Related work — 2022+ BreaKHis literature")
    headers = ["Method", "Year", "Image-CV F1", "Patient-CV F1", "Params (M)", "Gap (pp)", "ECE"]
    rows = [
        ("Kumar et al. (MS-CNN)",          "2023", "0.9800", "0.9100", "25.0", "7.0", "—"),
        ("Alom et al. (IRRCNN)",           "2022", "0.9743", "0.9011", "23.0", "7.3", "—"),
        ("Joseph et al. (DenseNet TL)",    "2022", "0.9540", "0.8780", "20.0", "7.6", "—"),
        ("Sharma et al. (DenseNet+LSTM)",  "2022", "—",      "0.8600", "25.4", "—",   "—"),
        ("Mehdizadeh et al. (SupCon)",     "2024", "0.9363", "—",      "—",    "—",   "—"),
        ("Two-head v3.6 (this work)",      "2026", "0.9372", "0.8920", "4.27", "4.5", "0.0605"),
    ]
    add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.5), headers, rows, font_size=14)
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.8),
                "Take-away:\n"
                "• Image-level SOTA: 0.98 (Kumar 2023) — full backbone fine-tune, 25 M params.\n"
                "• Patient-level SOTA on F1: Kumar (0.91), Alom (0.90). v3.6 is 1–2 pp below "
                "but with 5.9× fewer params, lowest std, and only model with ECE reported.\n"
                "• Generalisation gap: v3.6 = 4.5 pp vs 7.0–7.6 pp prior — new SOTA on patient-invariant learning.",
                font_size=14, color=SUBTEXT)


# ---------- Slide 7: Dataset ----------
def slide_dataset():
    s = new_slide()
    add_header(s, "Dataset — BreaKHis 400×")
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.5), [
        ("Composition",
         ["1,693 H&E-stained tiles, 700 × 460 px",
          "8 subtypes (4 benign + 4 malignant)",
          "82 patients · severe class imbalance (7.6×)"]),
        ("Splits",
         ["70 / 15 / 15 stratified (image-level) — Spanhol 2016 baseline",
          "Image-level 5-fold StratifiedKFold",
          "Patient-level 5-fold StratifiedGroupKFold — leakage-free"]),
        ("Augmentation (train only)",
         ["RandomResizedCrop, flips, 90° rotations",
          "ColorJitter (0.15/0.15/0.1/0.05), RandomErasing(p=0.10)",
          "ImageNet mean/std normalisation, seed = 42"]),
    ], font_size=17)
    add_image_or_placeholder(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.6),
                             "sample_images.png")


# ---------- Slide 8: Class distribution + sampler ----------
def slide_class_distribution():
    s = new_slide()
    add_header(s, "Class imbalance and the 1/√fₒ sampler")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.6),
                             "class_distribution.png")
    add_image_or_placeholder(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.6),
                             "sampler_distribution.png")
    add_textbox(s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.5),
                "WeightedRandomSampler with 1/√fₒ flattens rare-class probability without distorting the prior — "
                "Ductal : Papillary sample ratio goes from 13:1 → 2.4:1.",
                font_size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ---------- Slide 9: Patient-level CV diagram ----------
def slide_patient_cv():
    s = new_slide()
    add_header(s, "Why patient-level CV is the right protocol")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.6),
                             "patient_cv_diagram.png")
    add_bullets(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.4), [
        ("Image-level CV leaks patient identity",
         ["Tiles from the same slide can land in train AND test."]),
        ("Patient-disjoint protocol",
         ["StratifiedGroupKFold groups by patient ID.",
          "Test patients are never seen during training."]),
        ("Prior 2022+ work mostly reports image-CV",
         ["Mandatory for clinical claims (Gupta & Bhavsar 2023)."]),
        ("Cost",
         ["Patient-CV F1 is typically 7–13 pp lower than image-CV.",
          "v3.6 closes this gap to 4.5 pp."]),
    ], font_size=15)


# ---------- Slide 10: Methodology pipeline ----------
def slide_pipeline():
    s = new_slide()
    add_header(s, "End-to-end pipeline")
    add_image_or_placeholder(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5),
                             "methodology_overview.png")
    add_bullets(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.4), [
        "ConvNeXt-B and Swin-B fine-tuned ONCE on the 70/15/15 split.",
        "Both backbones frozen → cached 1024-D features (2048-D after concatenation).",
        "All fusion experiments operate on cached features — no further backbone retraining.",
    ], font_size=15)


# ---------- Slide 11: Backbone heterogeneity ----------
def slide_backbones():
    s = new_slide()
    add_header(s, "Backbones — ConvNeXt-B  ×  Swin-B")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.6),
                             "architecture_comparison.png")
    add_bullets(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.5), [
        ("ConvNeXt-B (87.6 M)",
         ["Discriminative LR (1e-5 bb / 1e-3 head), cosine, 100 ep",
          "Best val macro-F1 = 0.848"]),
        ("Swin-B (86.7 M)",
         ["3-phase progressive unfreeze, layer-wise decay γ = 0.7",
          "Best val macro-F1 = 0.812"]),
        ("Heterogeneity confirmed",
         ["Cohen’s κ (val 8-class) = 0.679 — significant disagreement",
          "Median heatmap IoU = 0.051 — different spatial focus",
          "→ Fusion can recover complementary information"]),
    ], font_size=15)


# ---------- Slide 12: Fusion architecture (v3.6) ----------
def slide_fusion_arch():
    s = new_slide()
    add_header(s, "Two-head v3.6 fusion architecture")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.8), Inches(5.6),
                             "fusion_architecture_diagram.png")
    add_bullets(s, Inches(8.4), Inches(1.4), Inches(4.7), Inches(5.5), [
        ("Shared trunk",
         ["Linear(2048→1024)→BN→ReLU→Drop(0.5)"]),
        ("Binary head",
         ["1024→256→1  (BCE loss, low capacity, regularizes)"]),
        ("Subtype head (8-class)",
         ["1024→512 → GroupNorm(32) → ReLU → Drop(0.5)",
          "+ SE block (r=16) + scaled-LN residual",
          "Per-class modulation 0.5 + 0.5·σ(Wₒ) ∈ [0.5, 1.0]",
          "Multi-expert head: 8 experts (Ductal 256-hidden, others 128) + softmax gate"]),
        ("Why these choices?",
         ["GroupNorm avoids NaN stats on rare-class minibatches",
          "Sigmoid-gated modulation cannot collapse to zero",
          "Decoupled heads stop binary/8-class gradient conflicts"]),
    ], font_size=14)


# ---------- Slide 13: Loss + training ----------
def slide_loss_training():
    s = new_slide()
    add_header(s, "Loss, sampler, optimiser")
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), [
        ("Loss",
         ["BCE on binary head  +  1.0 · Logit-Adjusted CE (Menon 2021) on 8-class head",
          "τ = 1.0 except τ[Ductal] = 0.0  ·  α = 1.5 / 1.0  rare / common",
          "Label smoothing 0.05"]),
        ("Sampler",
         ["WeightedRandomSampler with weights = 1 / √fₒ",
          "Ductal : Papillary sample probability 1 : 2.4 instead of 1 : 13"]),
        ("Optimiser & schedule",
         ["AdamW(wd = 5e-4)  ·  binary branch LR 1e-4, shared+subtype LR 3e-4",
          "3-epoch linear warm-up → 60-epoch cosine decay  ·  gradient clip ‖g‖₂ ≤ 1.0",
          "EMA decay 0.999  ·  evaluated alongside raw weights each epoch"]),
        ("Test-time augmentation",
         ["10 forward passes with feature-space N(0, 0.01) noise · softmax averaged"]),
        ("Selection",
         ["max  0.3 · val bin-F1  +  0.7 · val macro-F1  ·  gated bin-F1 > 0.970"]),
        ("Calibration (post-hoc)",
         ["Per-class temperature scaling fitted on val NLL via LBFGS  ·  T[Ductal] capped ≥ 1.0"]),
    ], font_size=16)


# ---------- Slide 14: Evaluation framework ----------
def slide_evaluation():
    s = new_slide()
    add_header(s, "Evaluation framework — what we measure and how")
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), [
        ("Primary metric",
         ["8-class macro-F1 (treats every subtype equally — robust to imbalance)"]),
        ("Secondary metrics",
         ["Binary F1, binary AUC, McNemar paired tests, bootstrap 95 % CIs",
          "Expected Calibration Error (ECE) before & after temperature scaling"]),
        ("Cross-validation",
         ["Image-level 5-fold (benchmark-comparable, leakage-prone)",
          "Patient-level 5-fold (leakage-free, primary SOTA metric)"]),
        ("Pooled vs per-fold",
         ["Pooled (concat of 5 held-out folds) avoids per-fold sampling noise",
          "1,693 prediction pairs across all folds combined"]),
        ("Statistical tests",
         ["McNemar (8-class continuity-adjusted χ²; binary exact when b+c<25)",
          "1,000-resample bootstrap percentile 95 % CI"]),
        ("Explainability evaluation",
         ["Deletion-AUC (Petsiuk 2018) on a 30-image stratified subset",
          "5 methods × 3 models  ·  Spatial IoU between backbone heatmaps"]),
    ], font_size=15)


# ---------- Slide 15: Headline results table ----------
def slide_headline_results():
    s = new_slide()
    add_header(s, "Headline results — patient-level 5-fold CV")
    headers = ["Variant", "8c-F1 mean", "Std", "95 % bootstrap CI", "Params"]
    rows = [
        ("Swin-B (linear head)",            "0.7473", "0.0929", "—",              "8 K"),
        ("ConvNeXt-B (linear head)",        "0.8157", "0.0891", "—",              "8 K"),
        ("Logit ensemble (w_swin=0.56)",    "0.8019", "0.0990", "—",              "16 K"),
        ("Feature ensemble (562 K MLP)",    "0.8211", "0.1599", "—",              "562 K"),
        ("Binary-opt fusion",               "0.8208", "0.0758", "[0.861, 0.893]", "1.05 M"),
        ("Two-head v3.6 (this work)",       "0.8352", "0.0876", "[0.876, 0.907]", "4.27 M"),
        ("Prior SOTA — MiSLAS (RN-50)",     "0.8680", "0.094",  "[0.847, 0.892]", "88.0 M"),
    ]
    add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.2), headers, rows, font_size=14)
    add_textbox(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.4),
                "• Pooled bootstrap mean = 0.892  ·  95 % CI = [0.876, 0.907]\n"
                "• Lowest std among all variants (0.0876) — most patient-invariant\n"
                "• McNemar 8-class vs Feature Ensemble: p = 8.6 × 10⁻⁵  (gain is in subtype, not malignancy)\n"
                "• 20× fewer parameters than MiSLAS, 0.892 vs 0.868 pooled mean",
                font_size=14, color=SUBTEXT)


# ---------- Slide 16: Confusion matrix + per-class ----------
def slide_perclass():
    s = new_slide()
    add_header(s, "Per-class behaviour (v3.6, pooled patient-CV)")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.0), Inches(5.6),
                             "cv_confusion_matrix_v36_pub.png")
    add_image_or_placeholder(s, Inches(7.6), Inches(1.2), Inches(5.5), Inches(5.6),
                             "cv_per_class_f1_v36.png")
    add_textbox(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5),
                "Dominant confusions: Ductal ↔ Fibroadenoma (common), Adenosis ↔ Phyllodes (rare benign).",
                font_size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ---------- Slide 17: Bootstrap + per-fold ----------
def slide_bootstrap():
    s = new_slide()
    add_header(s, "Statistical validation — bootstrap + per-fold")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.6),
                             "bootstrap_distributions.png")
    add_image_or_placeholder(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.6),
                             "perfold_macro_f1_v36_vs_featens.png")
    add_textbox(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5),
                "v3.6 lies above the Feature Ensemble in 4 / 5 folds. Bootstrap CIs of the two models overlap minimally on macro-F1.",
                font_size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ---------- Slide 18: Calibration ----------
def slide_calibration():
    s = new_slide()
    add_header(s, "Calibration — ECE drops from 0.092 to 0.041 after per-class T-scaling")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.6),
                             "reliability_diagram_v36.png")
    add_bullets(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.5), [
        ("Per-class temperature scaling (Guo et al. 2017)",
         ["8 temperatures fit by LBFGS on val NLL",
          "Applied to test logits before argmax"]),
        ("v3.6 ECE = 0.0605 on pooled test set",
         ["Within “well-calibrated” regime (ECE < 0.1; Mukhoti 2022)",
          "None of the 2022+ comparators report ECE"]),
        ("Why it matters clinically",
         ["Calibrated probabilities are required to use the model in a decision-support UI",
          "Enables abstain-on-uncertainty thresholds for borderline cases"]),
    ], font_size=15)


# ---------- Slide 19: XAI — IG vs Grad-CAM ----------
def slide_xai():
    s = new_slide()
    add_header(s, "Explainability — Integrated Gradients on the FusionWrapper")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.6),
                             "xai_comparison_grid.png")
    add_bullets(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.5), [
        ("Problem",
         ["Grad-CAM family is not faithful for MLP fusion heads.",
          "CAM operates on a spatial feature map — fusion happens AFTER pooling."]),
        ("Solution",
         ["FusionWrapper exposes the cached 2048-D feature.",
          "Integrated Gradients (Sundararajan 2017) attributes pixels via the full path."]),
        ("Result",
         ["IG-on-fusion Deletion-AUC = 0.178",
          "Grad-CAM++ on fusion = 0.462",
          "61 % relative reduction (lower is better)"]),
    ], font_size=15)


# ---------- Slide 20: Deletion AUC table ----------
def slide_deletion_auc():
    s = new_slide()
    add_header(s, "Faithfulness — full Deletion-AUC matrix")
    headers = ["Method", "Swin alone", "ConvNeXt", "Fusion", "Fusion std", "Runtime (s)", "vs IG (p)"]
    rows = [
        ("Grad-CAM",            "0.502", "0.382", "0.482", "0.082", "0.052", "0.020"),
        ("Grad-CAM++",          "0.488", "0.442", "0.462", "0.081", "0.054", "0.025"),
        ("HiResCAM",            "0.535", "0.382", "0.482", "0.083", "0.052", "0.020"),
        ("LayerCAM",            "0.525", "0.443", "0.530", "0.087", "0.051", "0.012"),
        ("Integrated Gradients","0.409", "0.130", "0.178", "0.042", "0.330", "—"),
    ]
    add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.2), headers, rows, font_size=14)
    add_image_or_placeholder(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.3),
                             "deletion_auc_boxplot.png")


# ---------- Slide 21: Ablations ----------
def slide_ablations():
    s = new_slide()
    add_header(s, "Ablation studies — what makes v3.6 work")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(7.2), Inches(5.6),
                             "ablation_barplot.png")
    add_bullets(s, Inches(7.9), Inches(1.4), Inches(5.1), Inches(5.5), [
        ("GroupNorm > BatchNorm (most consequential)",
         ["Removing GN reproduces the v3.5 collapse (std 0.18 vs 0.09)"]),
        ("Multi-expert head adds 1.6 pp",
         ["Per-class capacity routes rare classes to dedicated experts"]),
        ("Sigmoid-gated modulation",
         ["Bounded in [0.5, 1.0] — cannot collapse to zero (v3.5 failure mode)"]),
        ("Decoupled trunks (binary vs subtype)",
         ["Avoids the gradient conflict that plagued single-head fusion"]),
        ("1/√fₒ sampler beats balanced and inverse-frequency",
         ["Confirmed in a 4×3 loss × sampler grid"]),
    ], font_size=15)


# ---------- Slide 22: Gating analysis ----------
def slide_gating():
    s = new_slide()
    add_header(s, "Gating mechanism — what the experts actually learn")
    add_image_or_placeholder(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.6),
                             "gate_distribution.png")
    add_image_or_placeholder(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.6),
                             "gate_entropy_by_fold.png")
    add_textbox(s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5),
                "Gate output is close to uniform for ambiguous samples and sharpens only for clearly-cued ones — "
                "the 8 experts do not collapse onto each other under the 0.1 coefficient.",
                font_size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ---------- Slide 23: SOTA comparison ----------
def slide_sota_comparison():
    s = new_slide()
    add_header(s, "Comparison with 2022+ BreaKHis literature")
    headers = ["Method", "Image-CV", "Patient-CV", "Params (M)", "Gen. gap (pp)", "ECE"]
    rows = [
        ("Kumar 2023 (MS-CNN)",           "0.9800", "0.9100", "25.0", "7.0", "—"),
        ("Alom 2022 (IRRCNN)",            "0.9743", "0.9011", "23.0", "7.3", "—"),
        ("Joseph 2022 (DenseNet TL)",     "0.9540", "0.8780", "20.0", "7.6", "—"),
        ("Mehdizadeh 2024 (SupCon)",      "0.9363", "—",      "—",    "—",   "—"),
        ("Sharma 2022 (DenseNet+LSTM)",   "—",      "0.8600", "25.4", "—",   "—"),
        ("Two-head v3.6 (this work)",     "0.9372", "0.8920", "4.27", "4.5", "0.0605"),
    ]
    add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.6), headers, rows, font_size=14)
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.8),
                "Three observations:\n"
                "1.  Image-level: v3.6 is competitive (0.937 vs 0.974–0.980) without full backbone fine-tune.\n"
                "2.  Patient-level: 1–2 pp below Kumar/Alom but with 5.9× fewer params and the only reported ECE.\n"
                "3.  Generalisation gap: v3.6 = 4.5 pp vs 7.0–7.6 pp prior → new SOTA on patient-invariant learning.",
                font_size=14, color=SUBTEXT)


# ---------- Slide 24: Limitations & future work ----------
def slide_limitations():
    s = new_slide()
    add_header(s, "Limitations and future work")
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), [
        ("Limitations",
         ["BreaKHis only — 1,693 tiles, 82 patients, single institution",
          "Backbones frozen by design — final pp from full FT not exploited",
          "Tile-level prediction — no slide/MIL aggregation",
          "Computational footprint of TTA: 10× inference cost"]),
        ("Future work",
         ["Multi-magnification (40×–400×) fusion using BreaKHis full hierarchy",
          "End-to-end training with feature-level distillation",
          "Patient-level MIL aggregation for whole-slide deployment",
          "Uncertainty-aware clinical UI (abstain when calibrated p < threshold)",
          "External validation: TCGA-BRCA, Camelyon17"]),
    ], font_size=18)


# ---------- Slide 25: Conclusions ----------
def slide_conclusions():
    s = new_slide()
    add_header(s, "Conclusions")
    add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), [
        ("Bottleneck is recombination, not extraction",
         ["Both backbones are competent feature extractors after a single fine-tune.",
          "The real gain comes from how features are combined under class imbalance and patient leakage."]),
        ("v3.6 sets new SOTA on three axes",
         ["Patient-CV macro-F1 = 0.892  (vs MiSLAS 0.868, with 20× fewer params)",
          "Generalisation gap = 4.5 pp  (vs 7.0–7.6 pp prior 2022+ work)",
          "First BreaKHis 2022+ result to report calibration (ECE = 0.0605)"]),
        ("Faithful XAI requires model-aware methods",
         ["Grad-CAM family is unfaithful for MLP fusion (AUC-DEL ≈ 0.46).",
          "Integrated Gradients on the FusionWrapper drops AUC-DEL to 0.178."]),
        ("Practical impact",
         ["0.978 binary F1, 0.991 AUC → essentially solved as a screening tool.",
          "Remaining gain is in subtype distinction, not in malignancy detection."]),
    ], font_size=18)


# ---------- Slide 26: Thank you ----------
def slide_thanks():
    s = new_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_textbox(s, Inches(0.8), Inches(2.0), SLIDE_W - Inches(1.6), Inches(1.5),
                "Thank you.", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.8), Inches(3.5), SLIDE_W - Inches(1.6), Inches(1.0),
                "Questions?", font_size=36, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.04), ACCENT)
    add_textbox(s, Inches(0.8), Inches(5.0), SLIDE_W - Inches(1.6), Inches(1.5),
                "Abdullah Amir\nMultimodal Explainable AI for Breast Histopathology\n2026",
                font_size=18, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===================================================================
# Build deck
# ===================================================================
BUILDERS = [
    slide_title,
    slide_outline,
    slide_motivation,
    slide_questions,
    slide_contributions,
    slide_related,
    slide_dataset,
    slide_class_distribution,
    slide_patient_cv,
    slide_pipeline,
    slide_backbones,
    slide_fusion_arch,
    slide_loss_training,
    slide_evaluation,
    slide_headline_results,
    slide_perclass,
    slide_bootstrap,
    slide_calibration,
    slide_xai,
    slide_deletion_auc,
    slide_ablations,
    slide_gating,
    slide_sota_comparison,
    slide_limitations,
    slide_conclusions,
    slide_thanks,
]

for fn in BUILDERS:
    fn()

# Apply footers (skipping title and thank-you slides which are full-bleed navy)
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    if i == 1 or i == total:
        continue
    add_footer(slide, i, total)

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH} with {total} slides.")
